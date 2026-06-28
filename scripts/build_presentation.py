"""
build_presentation.py
======================

Auto-generate a 10-slide PowerPoint presentation for the BAH 2026 hackathon
from PRESENTATION_CONTENT.md. Run this to get a polished .pptx file.

Usage:
    D:\\BAH2026\\.venv\\Scripts\\python.exe D:\\BAH2026\\cross_modal_retrieval\\scripts\\build_presentation.py
    D:\\BAH2026\\.venv\\Scripts\\python.exe D:\\BAH2026\\cross_modal_retrieval\\scripts\\build_presentation.py --output D:\\BAH2026_presentation.pptx

After running, just open the .pptx in PowerPoint and tweak colors/images.
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Install python-pptx: pip install python-pptx")
    sys.exit(1)


# Theme colors (ISRO-style blue palette)
COL_BG = RGBColor(0x0A, 0x1F, 0x44)         # deep navy
COL_PRIMARY = RGBColor(0x4A, 0x90, 0xE2)    # sky blue
COL_ACCENT = RGBColor(0xFF, 0xB3, 0x00)     # ISRO orange
COL_TEXT = RGBColor(0xFF, 0xFF, 0xFF)       # white
COL_MUTED = RGBColor(0xB0, 0xC4, 0xDE)      # light blue-gray


def add_title_slide(prs, title, subtitle, author):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _bg(slide, COL_BG)
    _rect(slide, 0, Inches(6.5), prs.slide_width, Inches(1.5), COL_PRIMARY)
    _text(slide, title, Inches(0.5), Inches(1.5), prs.slide_width - Inches(1.0),
          Inches(2.5), font_size=40, bold=True, color=COL_TEXT, align=PP_ALIGN.CENTER)
    _text(slide, subtitle, Inches(0.5), Inches(4.2), prs.slide_width - Inches(1.0),
          Inches(1.0), font_size=22, color=COL_MUTED, align=PP_ALIGN.CENTER)
    _text(slide, author, Inches(0.5), Inches(6.7), prs.slide_width - Inches(1.0),
          Inches(0.5), font_size=14, color=COL_TEXT, align=PP_ALIGN.CENTER)


def add_content_slide(prs, title, bullets, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, COL_BG)
    _text(slide, title, Inches(0.5), Inches(0.3), prs.slide_width - Inches(1.0),
          Inches(0.9), font_size=32, bold=True, color=COL_ACCENT)
    _rect(slide, 0, Inches(1.2), prs.slide_width, Inches(0.04), COL_PRIMARY)
    body_y = Inches(1.5)
    for b in bullets:
        if isinstance(b, tuple) and b[0] == "TABLE":
            # b = ("TABLE", headers, rows, col_widths)
            _, headers, rows, _ = b
            _add_table(slide, headers, rows, Inches(0.6), body_y, prs.slide_width - Inches(1.2))
            body_y += Inches(0.5) + Inches(0.4 * (len(rows) + 1))
        else:
            _text(slide, "• " + b, Inches(0.7), body_y, prs.slide_width - Inches(1.4),
                  Inches(0.55), font_size=18, color=COL_TEXT)
            body_y += Inches(0.6)
    if footer:
        _text(slide, footer, Inches(0.5), prs.slide_height - Inches(0.4),
              prs.slide_width - Inches(1.0), Inches(0.3),
              font_size=10, color=COL_MUTED, align=PP_ALIGN.RIGHT)


# ---- helpers ----

def _bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  slide.part.package.presentation_part.presentation.slide_width,
                                  slide.part.package.presentation_part.presentation.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def _rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()


def _text(slide, text, x, y, w, h, font_size=18, bold=False, color=COL_TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color


def _add_table(slide, headers, rows, x, y, w):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, Inches(0.4 * n_rows))
    table = table_shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = COL_TEXT
        cell.fill.solid(); cell.fill.fore_color.rgb = COL_PRIMARY
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12); r.font.color.rgb = COL_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = COL_BG if ri % 2 == 0 else RGBColor(0x14, 0x2A, 0x55)


def build(out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(prs,
                    "Cross-Modal Satellite Image Retrieval",
                    "Using Multi-Sensor Remote Sensing Data · BAH 2026",
                    "Team Swaraj Kamila · ISRO BAH 2026")

    # 2. Problem
    add_content_slide(prs, "Problem Statement #11",
        ["Same geographic location captured by different sensors (Optical, MS, SAR)",
         "Metadata-based search fails across sensor types and seasons",
         "Goal: One query → top-K similar images, irrespective of sensor modality"],
        footer="BAH 2026 — ISRO Hackathon")

    # 3. Objectives
    add_content_slide(prs, "Objectives",
        ["Same-modal: optical→optical, SAR→SAR, MS→MS",
         "Cross-modal: optical↔SAR, optical↔MS, MS↔SAR",
         "Top-5 and Top-10 ranked retrieval lists",
         "Sub-millisecond retrieval latency",
         "F1@5, F1@10, P@k, Hit@k, MAP@k reported"],
        footer="Disaster response · Agriculture monitoring · Defense intel")

    # 4. Dataset
    add_content_slide(prs, "Multi-Sensor Dataset", [
        ("TABLE",
         ["Source", "Modality", "Format", "Classes", "Samples"],
         [["EuroSAT", "Multispectral (s1)", ".tif", "10", "~30K"],
          ["EuroSAT", "Optical (s2)", ".jpg", "10", "~30K"],
          ["Sentinel", "SAR (s1)", ".png", "4", "~2K"],
          ["Sentinel", "Optical (s2)", ".png", "4", "~2K"],
          ["Gallery", "3 modalities", "mixed", "14", "~9.2K"]],
         None),
        "Pairing: EuroSAT by image id; Sentinel by patch id (_p10.png)"],
        footer="Total gallery: ~9,200 images across 14 semantic classes")

    # 5. Architecture
    add_content_slide(prs, "End-to-End Pipeline",
        ["Raw Images (SAR / MS / Optical)",
         "   ↓",
         "DINOv2-Base @ 518px  →  768-D features (foundation model)",
         "   ↓",
         "Per-modality Projector V6 (1024 hidden, 2 residual blocks)",
         "   ↓",
         "256-D shared embedding (L2-normalized, cosine space)",
         "   ↓",
         "FAISS IndexFlatIP  →  ~10 MB on disk",
         "   ↓",
         "Top-K retrieval  ~0.02 ms per query"],
        footer="Stack: PyTorch · timm · FAISS · Flask")

    # 6. Methodology
    add_content_slide(prs, "Methodology",
        ["Stage 1 — Feature Extraction:",
         "   DINOv2-Base (LVD-142M) · 4-view TTA (orig + hflip + vflip + hvflip)",
         "Stage 2 — Shared Embedding Training:",
         "   2-block residual MLP · 50 epochs · AdamW + cosine LR",
         "   Loss = 0.30 Triplet + 0.35 InfoNCE + 0.20 CrossModal + 0.15 Prototype",
         "   Class-balanced sampling · hard negative mining",
         "Stage 3 — FAISS IndexFlatIP on L2-normalized 256-D vectors",
         "Stage 4 — Eval: 60 queries/class · F1@5/10 · P@k · Hit@k · MAP@k"],
        footer="")

    # 7. Innovations
    add_content_slide(prs, "What Makes This Work",
        ["Foundation model transfer: DINOv2 web-trained → strong semantic priors",
         "   Outperforms ResNet-50 on remote sensing without any fine-tuning",
         "Test-Time Augmentation: 4-view averaging reduces sensor noise",
         "   +0.5-1% precision for 4x compute",
         "Multi-loss training: Triplet + InfoNCE + CrossModal + Prototype",
         "   Stable gradients, mod-invariant embedding space",
         "Class-balanced sampling: every batch sees all 14 classes",
         "   Avoids Forest (3000) vs urban (800) bias"],
        footer="")

    # 8. Same-modal results
    add_content_slide(prs, "Results — Same-Modal Retrieval", [
        ("TABLE",
         ["Query → Gallery", "n", "P@5", "P@10", "Hit@5", "MAP@10", "Time (ms)"],
         [["MS → MS", "300", "0.992", "0.993", "1.000", "0.985", "0.027"],
          ["Optical → Optical", "420", "0.991", "0.992", "1.000", "0.984", "0.019"],
          ["SAR → SAR", "120", "1.000", "1.000", "1.000", "1.000", "0.012"],
          ["Average", "—", "0.994", "0.995", "1.000", "0.990", "0.019"]],
         None),
        "Hit rate = 100% — every query retrieves ≥1 relevant item"],
        footer="Sub-millisecond retrieval on 9.2K gallery")

    # 9. Cross-modal results
    add_content_slide(prs, "Results — Cross-Modal Retrieval", [
        ("TABLE",
         ["Query → Gallery", "n", "P@5", "P@10", "Hit@5", "MAP@10", "Time (ms)"],
         [["MS → Optical", "300", "0.991", "0.992", "1.000", "0.982", "0.020"],
          ["Optical → MS", "420", "0.989", "0.991", "1.000", "0.981", "0.020"],
          ["SAR → Optical", "120", "1.000", "1.000", "1.000", "1.000", "0.025"],
          ["Optical → SAR", "420", "0.993", "0.994", "1.000", "0.987", "0.022"],
          ["MS → SAR", "300", "0.992", "0.993", "1.000", "0.985", "0.013"],
          ["SAR → MS", "120", "1.000", "1.000", "1.000", "1.000", "0.012"],
          ["Average", "—", "0.994", "0.995", "1.000", "0.989", "0.019"]],
         None),
        "Cross-modal matches same-modal — proves modality-invariance"],
        footer="")

    # 10. Demo + conclusion
    add_content_slide(prs, "Demo & Live System",
        ["🌐 Live web app deployed on Render",
         "📂 Upload query from any modality (.tif / .png / .jpg)",
         "⚡ Top-5/Top-10 returned in <50 ms (extraction + search)",
         "🎯 Per-result: modality badge, similarity score, latency",
         "   Demo 1: Upload Optical → top includes MS + SAR of same class",
         "   Demo 2: Upload SAR → top includes Optical + MS of same land-cover"],
        footer="Live URL: cross-modal-retrieval.onrender.com")

    # 11. Conclusion + future work
    add_content_slide(prs, "Conclusion & Future Work",
        ["✅ 99%+ Precision@5 across all 9 (src, tgt) modality pairs",
         "✅ Modality-agnostic embedding space (DINOv2 + V6 projector)",
         "✅ Sub-millisecond retrieval at 9.2K gallery size",
         "✅ End-to-end: extract → train → index → serve → demo",
         "Future work:",
         "   • Scale to 100K+ gallery (FAISS-IVF / HNSW)",
         "   • Add hyperspectral modality",
         "   • Geo-coordinate-aware retrieval (not just class label)",
         "   • Active learning loop with user feedback"],
        footer="Thank you — questions welcome")

    prs.save(out_path)
    print(f"[pptx] saved -> {out_path}")
    print(f"[pptx] {len(prs.slides)} slides")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--output", type=Path, default=Path(r"D:\BAH2026_presentation.pptx"))
    args = ap.parse_args()
    build(args.output)


if __name__ == "__main__":
    main()